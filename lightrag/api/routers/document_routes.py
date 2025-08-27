"""
This module contains all document-related routes for the LightRAG API.
"""

import asyncio
from pyuca import Collator
from lightrag.utils import logger
import aiofiles
import shutil
import traceback
import pipmaster as pm
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Any, Literal
from fastapi import APIRouter, BackgroundTasks, Depends, File, HTTPException, UploadFile
from pydantic import BaseModel, Field, field_validator

from lightrag import LightRAG
from lightrag.base import DocProcessingStatus, DocStatus
from lightrag.api.utils_api import get_combined_auth_dependency
from ..config import global_args


# Function to format datetime to ISO format string with timezone information
def format_datetime(dt: Any) -> Optional[str]:
    """Format datetime to ISO format string with timezone information

    Args:
        dt: Datetime object, string, or None

    Returns:
        ISO format string with timezone information, or None if input is None
    """
    if dt is None:
        return None
    if isinstance(dt, str):
        return dt

    # Check if datetime object has timezone information
    if isinstance(dt, datetime):
        # If datetime object has no timezone info (naive datetime), add UTC timezone
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=timezone.utc)

    # Return ISO format string with timezone information
    return dt.isoformat()


router = APIRouter(
    prefix="/documents",
    tags=["documents"],
)

# Temporary file prefix
temp_prefix = "__tmp__"


class ScanResponse(BaseModel):
    """Response model for document scanning operation

    Attributes:
        status: Status of the scanning operation
        message: Optional message with additional details
    """

    status: Literal["scanning_started"] = Field(
        description="Status of the scanning operation"
    )
    message: Optional[str] = Field(
        default=None, description="Additional details about the scanning operation"
    )

    class Config:
        json_schema_extra = {
            "example": {
                "status": "scanning_started",
                "message": "Scanning process has been initiated in the background",
            }
        }


class InsertTextRequest(BaseModel):
    """Request model for inserting a single text document

    Attributes:
        text: The text content to be inserted into the RAG system
        file_source: Source of the text (optional)
    """

    text: str = Field(
        min_length=1,
        description="The text to insert",
    )
    file_source: str = Field(default=None, min_length=0, description="File Source")

    @field_validator("text", mode="after")
    @classmethod
    def strip_text_after(cls, text: str) -> str:
        return text.strip()

    @field_validator("file_source", mode="after")
    @classmethod
    def strip_source_after(cls, file_source: str) -> str:
        return file_source.strip()

    class Config:
        json_schema_extra = {
            "example": {
                "text": "This is a sample text to be inserted into the RAG system.",
                "file_source": "Source of the text (optional)",
            }
        }


class InsertTextsRequest(BaseModel):
    """Request model for inserting multiple text documents

    Attributes:
        texts: List of text contents to be inserted into the RAG system
        file_sources: Sources of the texts (optional)
    """

    texts: list[str] = Field(
        min_length=1,
        description="The texts to insert",
    )
    file_sources: list[str] = Field(
        default=None, min_length=0, description="Sources of the texts"
    )

    @field_validator("texts", mode="after")
    @classmethod
    def strip_texts_after(cls, texts: list[str]) -> list[str]:
        return [text.strip() for text in texts]

    @field_validator("file_sources", mode="after")
    @classmethod
    def strip_sources_after(cls, file_sources: list[str]) -> list[str]:
        return [file_source.strip() for file_source in file_sources]

    class Config:
        json_schema_extra = {
            "example": {
                "texts": [
                    "This is the first text to be inserted.",
                    "This is the second text to be inserted.",
                ],
                "file_sources": [
                    "First file source (optional)",
                ],
            }
        }


class N8nWebhookRequest(BaseModel):
    """Request model for n8n webhook integration

    Attributes:
        content: The main content to be inserted (can be text or JSON string)
        source: Source identifier for the content (optional)
        metadata: Additional metadata from n8n (optional)
        content_type: Type of content (text, json, file) (optional)
    """

    content: str = Field(
        min_length=1,
        description="The content to insert (text, JSON string, or file content)",
    )
    source: Optional[str] = Field(
        default="n8n_webhook", 
        description="Source identifier for the content"
    )
    metadata: Optional[Dict[str, Any]] = Field(
        default=None,
        description="Additional metadata from n8n workflow"
    )
    content_type: Optional[str] = Field(
        default="text",
        description="Type of content: 'text', 'json', or 'file'"
    )

    @field_validator("content", mode="after")
    @classmethod
    def strip_content_after(cls, content: str) -> str:
        return content.strip()

    @field_validator("source", mode="after")
    @classmethod
    def strip_source_after(cls, source: str) -> str:
        return source.strip() if source else "n8n_webhook"

    class Config:
        json_schema_extra = {
            "example": {
                "content": "This is sample content from n8n workflow",
                "source": "email_processor",
                "metadata": {
                    "workflow_id": "12345",
                    "trigger_type": "email",
                    "timestamp": "2024-01-01T00:00:00Z"
                },
                "content_type": "text"
            }
        }


class InsertResponse(BaseModel):
    """Response model for document insertion operations

    Attributes:
        status: Status of the operation (success, duplicated, partial_success, failure)
        message: Detailed message describing the operation result
    """

    status: Literal["success", "duplicated", "partial_success", "failure"] = Field(
        description="Status of the operation"
    )
    message: str = Field(description="Message describing the operation result")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "success",
                "message": "File 'document.pdf' uploaded successfully. Processing will continue in background.",
            }
        }


class ClearDocumentsResponse(BaseModel):
    """Response model for document clearing operation

    Attributes:
        status: Status of the clear operation
        message: Detailed message describing the operation result
    """

    status: Literal["success", "partial_success", "busy", "fail"] = Field(
        description="Status of the clear operation"
    )
    message: str = Field(description="Message describing the operation result")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "success",
                "message": "All documents cleared successfully. Deleted 15 files.",
            }
        }


class ClearCacheRequest(BaseModel):
    """Request model for clearing cache

    Attributes:
        modes: Optional list of cache modes to clear
    """

    modes: Optional[
        List[Literal["default", "naive", "local", "global", "hybrid", "mix"]]
    ] = Field(
        default=None,
        description="Modes of cache to clear. If None, clears all cache.",
    )

    class Config:
        json_schema_extra = {"example": {"modes": ["default", "naive"]}}


class ClearCacheResponse(BaseModel):
    """Response model for cache clearing operation

    Attributes:
        status: Status of the clear operation
        message: Detailed message describing the operation result
    """

    status: Literal["success", "fail"] = Field(
        description="Status of the clear operation"
    )
    message: str = Field(description="Message describing the operation result")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "success",
                "message": "Successfully cleared cache for modes: ['default', 'naive']",
            }
        }


"""Response model for document status

Attributes:
    id: Document identifier
    content_summary: Summary of document content
    content_length: Length of document content
    status: Current processing status
    created_at: Creation timestamp (ISO format string)
    updated_at: Last update timestamp (ISO format string)
    chunks_count: Number of chunks (optional)
    error: Error message if any (optional)
    metadata: Additional metadata (optional)
    file_path: Path to the document file
"""


class DocStatusResponse(BaseModel):
    id: str = Field(description="Document identifier")
    content_summary: str = Field(description="Summary of document content")
    content_length: int = Field(description="Length of document content in characters")
    status: DocStatus = Field(description="Current processing status")
    created_at: str = Field(description="Creation timestamp (ISO format string)")
    updated_at: str = Field(description="Last update timestamp (ISO format string)")
    chunks_count: Optional[int] = Field(
        default=None, description="Number of chunks the document was split into"
    )
    error: Optional[str] = Field(
        default=None, description="Error message if processing failed"
    )
    metadata: Optional[dict[str, Any]] = Field(
        default=None, description="Additional metadata about the document"
    )
    file_path: str = Field(description="Path to the document file")

    class Config:
        json_schema_extra = {
            "example": {
                "id": "doc_123456",
                "content_summary": "Research paper on machine learning",
                "content_length": 15240,
                "status": "PROCESSED",
                "created_at": "2025-03-31T12:34:56",
                "updated_at": "2025-03-31T12:35:30",
                "chunks_count": 12,
                "error": None,
                "metadata": {"author": "John Doe", "year": 2025},
                "file_path": "research_paper.pdf",
            }
        }


class DocsStatusesResponse(BaseModel):
    """Response model for document statuses

    Attributes:
        statuses: Dictionary mapping document status to lists of document status responses
    """

    statuses: Dict[DocStatus, List[DocStatusResponse]] = Field(
        default_factory=dict,
        description="Dictionary mapping document status to lists of document status responses",
    )

    class Config:
        json_schema_extra = {
            "example": {
                "statuses": {
                    "PENDING": [
                        {
                            "id": "doc_123",
                            "content_summary": "Pending document",
                            "content_length": 5000,
                            "status": "PENDING",
                            "created_at": "2025-03-31T10:00:00",
                            "updated_at": "2025-03-31T10:00:00",
                            "file_path": "pending_doc.pdf",
                        }
                    ],
                    "PROCESSED": [
                        {
                            "id": "doc_456",
                            "content_summary": "Processed document",
                            "content_length": 8000,
                            "status": "PROCESSED",
                            "created_at": "2025-03-31T09:00:00",
                            "updated_at": "2025-03-31T09:05:00",
                            "chunks_count": 8,
                            "file_path": "processed_doc.pdf",
                        }
                    ],
                }
            }
        }


class PipelineStatusResponse(BaseModel):
    """Response model for pipeline status

    Attributes:
        autoscanned: Whether auto-scan has started
        busy: Whether the pipeline is currently busy
        job_name: Current job name (e.g., indexing files/indexing texts)
        job_start: Job start time as ISO format string with timezone (optional)
        docs: Total number of documents to be indexed
        batchs: Number of batches for processing documents
        cur_batch: Current processing batch
        request_pending: Flag for pending request for processing
        latest_message: Latest message from pipeline processing
        history_messages: List of history messages
        update_status: Status of update flags for all namespaces
    """

    autoscanned: bool = False
    busy: bool = False
    job_name: str = "Default Job"
    job_start: Optional[str] = None
    docs: int = 0
    batchs: int = 0
    cur_batch: int = 0
    request_pending: bool = False
    latest_message: str = ""
    history_messages: Optional[List[str]] = None
    update_status: Optional[dict] = None

    @field_validator("job_start", mode="before")
    @classmethod
    def parse_job_start(cls, value):
        """Process datetime and return as ISO format string with timezone"""
        return format_datetime(value)

    class Config:
        extra = "allow"  # Allow additional fields from the pipeline status


class DocumentManager:
    def __init__(
        self,
        input_dir: str,
        supported_extensions: tuple = (
            ".txt",
            ".md",
            ".pdf",
            ".docx",
            ".pptx",
            ".xlsx",
            ".rtf",  # Rich Text Format
            ".odt",  # OpenDocument Text
            ".tex",  # LaTeX
            ".epub",  # Electronic Publication
            ".html",  # HyperText Markup Language
            ".htm",  # HyperText Markup Language
            ".csv",  # Comma-Separated Values
            ".json",  # JavaScript Object Notation
            ".xml",  # eXtensible Markup Language
            ".yaml",  # YAML Ain't Markup Language
            ".yml",  # YAML
            ".log",  # Log files
            ".conf",  # Configuration files
            ".ini",  # Initialization files
            ".properties",  # Java properties files
            ".sql",  # SQL scripts
            ".bat",  # Batch files
            ".sh",  # Shell scripts
            ".c",  # C source code
            ".cpp",  # C++ source code
            ".py",  # Python source code
            ".java",  # Java source code
            ".js",  # JavaScript source code
            ".ts",  # TypeScript source code
            ".swift",  # Swift source code
            ".go",  # Go source code
            ".rb",  # Ruby source code
            ".php",  # PHP source code
            ".css",  # Cascading Style Sheets
            ".scss",  # Sassy CSS
            ".less",  # LESS CSS
        ),
    ):
        self.input_dir = Path(input_dir)
        self.supported_extensions = supported_extensions
        self.indexed_files = set()

        # Create input directory if it doesn't exist
        self.input_dir.mkdir(parents=True, exist_ok=True)

    def scan_directory_for_new_files(self) -> List[Path]:
        """Scan input directory for new files"""
        new_files = []
        for ext in self.supported_extensions:
            logger.debug(f"Scanning for {ext} files in {self.input_dir}")
            for file_path in self.input_dir.rglob(f"*{ext}"):
                if file_path not in self.indexed_files:
                    new_files.append(file_path)
        return new_files

    def mark_as_indexed(self, file_path: Path):
        self.indexed_files.add(file_path)

    def is_supported_file(self, filename: str) -> bool:
        return any(filename.lower().endswith(ext) for ext in self.supported_extensions)


async def pipeline_enqueue_file(rag: LightRAG, file_path: Path) -> bool:
    """Add a file to the queue for processing

    Args:
        rag: LightRAG instance
        file_path: Path to the saved file
    Returns:
        bool: True if the file was successfully enqueued, False otherwise
    """

    try:
        content = ""
        ext = file_path.suffix.lower()

        file = None
        async with aiofiles.open(file_path, "rb") as f:
            file = await f.read()

        # Process based on file type
        match ext:
            case (
                ".txt"
                | ".md"
                | ".html"
                | ".htm"
                | ".tex"
                | ".json"
                | ".xml"
                | ".yaml"
                | ".yml"
                | ".rtf"
                | ".odt"
                | ".epub"
                | ".csv"
                | ".log"
                | ".conf"
                | ".ini"
                | ".properties"
                | ".sql"
                | ".bat"
                | ".sh"
                | ".c"
                | ".cpp"
                | ".py"
                | ".java"
                | ".js"
                | ".ts"
                | ".swift"
                | ".go"
                | ".rb"
                | ".php"
                | ".css"
                | ".scss"
                | ".less"
            ):
                try:
                    # Try to decode as UTF-8
                    content = file.decode("utf-8")

                    # Validate content
                    if not content or len(content.strip()) == 0:
                        logger.error(f"Empty content in file: {file_path.name}")
                        return False

                    # Check if content looks like binary data string representation
                    if content.startswith("b'") or content.startswith('b"'):
                        logger.error(
                            f"File {file_path.name} appears to contain binary data representation instead of text"
                        )
                        return False

                except UnicodeDecodeError:
                    logger.error(
                        f"File {file_path.name} is not valid UTF-8 encoded text. Please convert it to UTF-8 before processing."
                    )
                    return False
            case ".pdf":
                if global_args.document_loading_engine == "DOCLING":
                    if not pm.is_installed("docling"):  # type: ignore
                        pm.install("docling")
                    from docling.document_converter import DocumentConverter  # type: ignore

                    converter = DocumentConverter()
                    result = converter.convert(file_path)
                    content = result.document.export_to_markdown()
                else:
                    if not pm.is_installed("pypdf2"):  # type: ignore
                        pm.install("pypdf2")
                    from PyPDF2 import PdfReader  # type: ignore
                    from io import BytesIO

                    pdf_file = BytesIO(file)
                    reader = PdfReader(pdf_file)
                    for page in reader.pages:
                        content += page.extract_text() + "\n"
            case ".docx":
                if global_args.document_loading_engine == "DOCLING":
                    if not pm.is_installed("docling"):  # type: ignore
                        pm.install("docling")
                    from docling.document_converter import DocumentConverter  # type: ignore

                    converter = DocumentConverter()
                    result = converter.convert(file_path)
                    content = result.document.export_to_markdown()
                else:
                    if not pm.is_installed("python-docx"):  # type: ignore
                        try:
                            pm.install("python-docx")
                        except Exception:
                            pm.install("docx")
                    from docx import Document  # type: ignore
                    from io import BytesIO

                    docx_file = BytesIO(file)
                    doc = Document(docx_file)
                    content = "\n".join(
                        [paragraph.text for paragraph in doc.paragraphs]
                    )
            case ".pptx":
                if global_args.document_loading_engine == "DOCLING":
                    if not pm.is_installed("docling"):  # type: ignore
                        pm.install("docling")
                    from docling.document_converter import DocumentConverter  # type: ignore

                    converter = DocumentConverter()
                    result = converter.convert(file_path)
                    content = result.document.export_to_markdown()
                else:
                    if not pm.is_installed("python-pptx"):  # type: ignore
                        pm.install("pptx")
                    from pptx import Presentation  # type: ignore
                    from io import BytesIO

                    pptx_file = BytesIO(file)
                    prs = Presentation(pptx_file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
            case ".xlsx":
                if global_args.document_loading_engine == "DOCLING":
                    if not pm.is_installed("docling"):  # type: ignore
                        pm.install("docling")
                    from docling.document_converter import DocumentConverter  # type: ignore

                    converter = DocumentConverter()
                    result = converter.convert(file_path)
                    content = result.document.export_to_markdown()
                else:
                    if not pm.is_installed("openpyxl"):  # type: ignore
                        pm.install("openpyxl")
                    from openpyxl import load_workbook  # type: ignore
                    from io import BytesIO

                    xlsx_file = BytesIO(file)
                    wb = load_workbook(xlsx_file)
                    for sheet in wb:
                        content += f"Sheet: {sheet.title}\n"
                        for row in sheet.iter_rows(values_only=True):
                            content += (
                                "\t".join(
                                    str(cell) if cell is not None else ""
                                    for cell in row
                                )
                                + "\n"
                            )
                        content += "\n"
            case _:
                logger.error(
                    f"Unsupported file type: {file_path.name} (extension {ext})"
                )
                return False

        # Insert into the RAG queue
        if content:
            await rag.apipeline_enqueue_documents(content, file_paths=file_path.name)
            logger.info(f"Successfully fetched and enqueued file: {file_path.name}")
            return True
        else:
            logger.error(f"No content could be extracted from file: {file_path.name}")

    except Exception as e:
        logger.error(f"Error processing or enqueueing file {file_path.name}: {str(e)}")
        logger.error(traceback.format_exc())
    finally:
        if file_path.name.startswith(temp_prefix):
            try:
                file_path.unlink()
            except Exception as e:
                logger.error(f"Error deleting file {file_path}: {str(e)}")
    return False


async def pipeline_index_file(rag: LightRAG, file_path: Path):
    """Index a file

    Args:
        rag: LightRAG instance
        file_path: Path to the saved file
    """
    try:
        if await pipeline_enqueue_file(rag, file_path):
            await rag.apipeline_process_enqueue_documents()

    except Exception as e:
        logger.error(f"Error indexing file {file_path.name}: {str(e)}")
        logger.error(traceback.format_exc())


async def pipeline_index_files(rag: LightRAG, file_paths: List[Path]):
    """Index multiple files sequentially to avoid high CPU load

    Args:
        rag: LightRAG instance
        file_paths: Paths to the files to index
    """
    if not file_paths:
        return
    try:
        enqueued = False

        # Create Collator for Unicode sorting
        collator = Collator()
        sorted_file_paths = sorted(file_paths, key=lambda p: collator.sort_key(str(p)))

        # Process files sequentially
        for file_path in sorted_file_paths:
            if await pipeline_enqueue_file(rag, file_path):
                enqueued = True

        # Process the queue only if at least one file was successfully enqueued
        if enqueued:
            await rag.apipeline_process_enqueue_documents()
    except Exception as e:
        logger.error(f"Error indexing files: {str(e)}")
        logger.error(traceback.format_exc())


async def pipeline_index_texts(
    rag: LightRAG, texts: List[str], file_sources: List[str] = None
):
    """Index a list of texts

    Args:
        rag: LightRAG instance
        texts: The texts to index
        file_sources: Sources of the texts
    """
    if not texts:
        return
    if file_sources is not None:
        if len(file_sources) != 0 and len(file_sources) != len(texts):
            [
                file_sources.append("unknown_source")
                for _ in range(len(file_sources), len(texts))
            ]
    await rag.apipeline_enqueue_documents(input=texts, file_paths=file_sources)
    await rag.apipeline_process_enqueue_documents()


# TODO: deprecate after /insert_file is removed
async def save_temp_file(input_dir: Path, file: UploadFile = File(...)) -> Path:
    """Save the uploaded file to a temporary location

    Args:
        file: The uploaded file

    Returns:
        Path: The path to the saved file
    """
    # Generate unique filename to avoid conflicts
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    unique_filename = f"{temp_prefix}{timestamp}_{file.filename}"

    # Create a temporary file to save the uploaded content
    temp_path = input_dir / "temp" / unique_filename
    temp_path.parent.mkdir(exist_ok=True)

    # Save the file
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    return temp_path


async def run_scanning_process(rag: LightRAG, doc_manager: DocumentManager):
    """Background task to scan and index documents"""
    try:
        new_files = doc_manager.scan_directory_for_new_files()
        total_files = len(new_files)
        logger.info(f"Found {total_files} new files to index.")

        if not new_files:
            return

        # Process all files at once
        await pipeline_index_files(rag, new_files)
        logger.info(f"Scanning process completed: {total_files} files Processed.")

    except Exception as e:
        logger.error(f"Error during scanning process: {str(e)}")
        logger.error(traceback.format_exc())


def create_document_routes(
    rag: LightRAG, doc_manager: DocumentManager, api_key: Optional[str] = None
):
    # Create combined auth dependency for document routes
    combined_auth = get_combined_auth_dependency(api_key)

    @router.post(
        "/scan", response_model=ScanResponse, dependencies=[Depends(combined_auth)]
    )
    async def scan_for_new_documents(background_tasks: BackgroundTasks):
        """
        Trigger the scanning process for new documents.

        This endpoint initiates a background task that scans the input directory for new documents
        and processes them. If a scanning process is already running, it returns a status indicating
        that fact.

        Returns:
            ScanResponse: A response object containing the scanning status
        """
        # Start the scanning process in the background
        background_tasks.add_task(run_scanning_process, rag, doc_manager)
        return ScanResponse(
            status="scanning_started",
            message="Scanning process has been initiated in the background",
        )

    @router.post(
        "/upload", response_model=InsertResponse, dependencies=[Depends(combined_auth)]
    )
    async def upload_to_input_dir(
        background_tasks: BackgroundTasks, file: UploadFile = File(...)
    ):
        """
        Upload a file to the input directory and index it.

        This API endpoint accepts a file through an HTTP POST request, checks if the
        uploaded file is of a supported type, saves it in the specified input directory,
        indexes it for retrieval, and returns a success status with relevant details.

        Args:
            background_tasks: FastAPI BackgroundTasks for async processing
            file (UploadFile): The file to be uploaded. It must have an allowed extension.

        Returns:
            InsertResponse: A response object containing the upload status and a message.
                status can be "success", "duplicated", or error is thrown.

        Raises:
            HTTPException: If the file type is not supported (400) or other errors occur (500).
        """
        try:
            if not doc_manager.is_supported_file(file.filename):
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type. Supported types: {doc_manager.supported_extensions}",
                )

            file_path = doc_manager.input_dir / file.filename
            # Check if file already exists
            if file_path.exists():
                return InsertResponse(
                    status="duplicated",
                    message=f"File '{file.filename}' already exists in the input directory.",
                )

            with open(file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)

            # Try to auto-link this upload to a pending document request by documentType
            try:
                from lightrag.api.routers.document_request_routes import (
                    document_requests_db,
                    get_current_timestamp,
                )

                # Infer documentType from filename prefix (e.g., "Tax Assessments.pdf" -> "tax_assessments")
                inferred_document_type = file.filename.rsplit(".", 1)[0].strip().lower().replace(" ", "_")

                # Find the most recent pending/requested request with matching documentType
                matching_id = None
                newest_ts = None
                for req_id, req in document_requests_db.items():
                    if req.get("status") in ("Requested", "Processing") and req.get("documentType"):
                        dt = str(req.get("documentType")).strip().lower().replace(" ", "_")
                        if dt == inferred_document_type:
                            created_ts = req.get("createdAt")
                            if newest_ts is None or (created_ts and created_ts > newest_ts):
                                newest_ts = created_ts
                                matching_id = req_id

                if matching_id:
                    # Load file content to store in request for download
                    with open(file_path, "rb") as fbin:
                        file_bytes = fbin.read()

                    document_requests_db[matching_id].update(
                        {
                            "status": "Ready",
                            "fileName": file.filename,
                            "fileSize": file_path.stat().st_size,
                            "updatedAt": get_current_timestamp(),
                            "downloadUrl": f"/webhook/api/document-requests/{matching_id}/download",
                            "file_content": file_bytes,
                        }
                    )
            except Exception as _auto_link_err:
                # Do not fail the upload if linking fails; log and continue
                logger.debug(f"Auto-link upload to pending request failed: {_auto_link_err}")

            # Add to background tasks
            background_tasks.add_task(pipeline_index_file, rag, file_path)

            return InsertResponse(
                status="success",
                message=f"File '{file.filename}' uploaded successfully. Processing will continue in background.",
            )
        except Exception as e:
            logger.error(f"Error /documents/upload: {file.filename}: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/text", response_model=InsertResponse, dependencies=[Depends(combined_auth)]
    )
    async def insert_text(
        request: InsertTextRequest, background_tasks: BackgroundTasks
    ):
        """
        Insert text into the RAG system.

        This endpoint allows you to insert text data into the RAG system for later retrieval
        and use in generating responses.

        Args:
            request (InsertTextRequest): The request body containing the text to be inserted.
            background_tasks: FastAPI BackgroundTasks for async processing

        Returns:
            InsertResponse: A response object containing the status of the operation.

        Raises:
            HTTPException: If an error occurs during text processing (500).
        """
        try:
            background_tasks.add_task(
                pipeline_index_texts,
                rag,
                [request.text],
                file_sources=[request.file_source],
            )
            return InsertResponse(
                status="success",
                message="Text successfully received. Processing will continue in background.",
            )
        except Exception as e:
            logger.error(f"Error /documents/text: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/texts",
        response_model=InsertResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def insert_texts(
        request: InsertTextsRequest, background_tasks: BackgroundTasks
    ):
        """
        Insert multiple texts into the RAG system.

        This endpoint allows you to insert multiple text entries into the RAG system
        in a single request.

        Args:
            request (InsertTextsRequest): The request body containing the list of texts.
            background_tasks: FastAPI BackgroundTasks for async processing

        Returns:
            InsertResponse: A response object containing the status of the operation.

        Raises:
            HTTPException: If an error occurs during text processing (500).
        """
        try:
            background_tasks.add_task(
                pipeline_index_texts,
                rag,
                request.texts,
                file_sources=request.file_sources,
            )
            return InsertResponse(
                status="success",
                message="Text successfully received. Processing will continue in background.",
            )
        except Exception as e:
            logger.error(f"Error /documents/text: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    # TODO: deprecated, use /upload instead
    @router.post(
        "/file", response_model=InsertResponse, dependencies=[Depends(combined_auth)]
    )
    async def insert_file(
        background_tasks: BackgroundTasks, file: UploadFile = File(...)
    ):
        """
        Insert a file directly into the RAG system.

        This endpoint accepts a file upload and processes it for inclusion in the RAG system.
        The file is saved temporarily and processed in the background.

        Args:
            background_tasks: FastAPI BackgroundTasks for async processing
            file (UploadFile): The file to be processed

        Returns:
            InsertResponse: A response object containing the status of the operation.

        Raises:
            HTTPException: If the file type is not supported (400) or other errors occur (500).
        """
        try:
            if not doc_manager.is_supported_file(file.filename):
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type. Supported types: {doc_manager.supported_extensions}",
                )

            temp_path = await save_temp_file(doc_manager.input_dir, file)

            # Add to background tasks
            background_tasks.add_task(pipeline_index_file, rag, temp_path)

            return InsertResponse(
                status="success",
                message=f"File '{file.filename}' saved successfully. Processing will continue in background.",
            )
        except Exception as e:
            logger.error(f"Error /documents/file: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    # TODO: deprecated, use /upload instead
    @router.post(
        "/file_batch",
        response_model=InsertResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def insert_batch(
        background_tasks: BackgroundTasks, files: List[UploadFile] = File(...)
    ):
        """
        Process multiple files in batch mode.

        This endpoint allows uploading and processing multiple files simultaneously.
        It handles partial successes and provides detailed feedback about failed files.

        Args:
            background_tasks: FastAPI BackgroundTasks for async processing
            files (List[UploadFile]): List of files to process

        Returns:
            InsertResponse: A response object containing:
                - status: "success", "partial_success", or "failure"
                - message: Detailed information about the operation results

        Raises:
            HTTPException: If an error occurs during processing (500).
        """
        try:
            inserted_count = 0
            failed_files = []
            temp_files = []

            for file in files:
                if doc_manager.is_supported_file(file.filename):
                    # Create a temporary file to save the uploaded content
                    temp_files.append(await save_temp_file(doc_manager.input_dir, file))
                    inserted_count += 1
                else:
                    failed_files.append(f"{file.filename} (unsupported type)")

            if temp_files:
                background_tasks.add_task(pipeline_index_files, rag, temp_files)

            # Prepare status message
            if inserted_count == len(files):
                status = "success"
                status_message = f"Successfully inserted all {inserted_count} documents"
            elif inserted_count > 0:
                status = "partial_success"
                status_message = f"Successfully inserted {inserted_count} out of {len(files)} documents"
                if failed_files:
                    status_message += f". Failed files: {', '.join(failed_files)}"
            else:
                status = "failure"
                status_message = "No documents were successfully inserted"
                if failed_files:
                    status_message += f". Failed files: {', '.join(failed_files)}"

            return InsertResponse(status=status, message=status_message)
        except Exception as e:
            logger.error(f"Error /documents/batch: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    # Conditionally expose n8n webhook based on feature flag
    if getattr(global_args, "enable_n8n_integration", False):
        @router.post(
            "/n8n_webhook",
            response_model=InsertResponse,
            dependencies=[Depends(combined_auth)],
        )
        async def n8n_webhook(
            request: N8nWebhookRequest, background_tasks: BackgroundTasks
        ):
        """
        Webhook endpoint for n8n integration.

        This endpoint is specifically designed for n8n workflows to send content
        to the LightRAG system. It accepts various content types and processes
        them in the background.

        Args:
            request (N8nWebhookRequest): The request body containing:
                - content: The main content to be inserted
                - source: Source identifier (optional, defaults to "n8n_webhook")
                - metadata: Additional metadata from n8n workflow (optional)
                - content_type: Type of content ("text", "json", "file") (optional)
            background_tasks: FastAPI BackgroundTasks for async processing

        Returns:
            InsertResponse: A response object containing the status of the operation.

        Raises:
            HTTPException: If an error occurs during processing (500).

        Example n8n webhook payload:
        {
            "content": "This is content from n8n workflow",
            "source": "email_processor",
            "metadata": {
                "workflow_id": "12345",
                "trigger_type": "email",
                "timestamp": "2024-01-01T00:00:00Z"
            },
            "content_type": "text"
        }
        """
            try:
            # Process content based on content_type
            if request.content_type == "json":
                # For JSON content, we might want to extract specific fields
                # or convert to a more readable format
                import json
                try:
                    json_data = json.loads(request.content)
                    # Convert JSON to a readable string format
                    processed_content = json.dumps(json_data, indent=2, ensure_ascii=False)
                except json.JSONDecodeError:
                    # If JSON parsing fails, treat as regular text
                    processed_content = request.content
            else:
                # For text and file content types, use as-is
                processed_content = request.content

            # Create a source identifier that includes metadata if available
            source_identifier = request.source
            if request.metadata:
                workflow_id = request.metadata.get("workflow_id", "unknown")
                source_identifier = f"{request.source}_{workflow_id}"

            # Create document request entry for frontend display
            try:
                from lightrag.api.routers.document_request_routes import get_current_timestamp
                from lightrag.kg.shared_storage import get_namespace_data, get_storage_lock
                import uuid
                
                request_id = str(uuid.uuid4())
                db = await get_namespace_data("document_requests")
                async with get_storage_lock():
                    db[request_id] = {
                        "requestId": request_id,
                        "status": "Processing",
                        "documentType": request.metadata.get("document_type", "n8n_upload") if request.metadata else "n8n_upload",
                        "parameters": request.metadata,
                        "downloadUrl": None,
                        "fileName": f"{source_identifier}.{request.content_type}",
                        "fileSize": len(request.content),
                        "errorMessage": None,
                        "createdAt": get_current_timestamp(),
                        "updatedAt": get_current_timestamp(),
                        "file_content": request.content.encode('utf-8'),  # Store actual content for download
                        "lightrag_source_id": source_identifier  # Link to LightRAG document
                    }
                
                # Add background task to monitor LightRAG processing
                async def monitor_lightrag_processing():
                    import asyncio
                    max_attempts = 60  # 5 minutes max (60 * 5 seconds)
                    attempts = 0
                    
                    while attempts < max_attempts:
                        try:
                            # Check if document exists in LightRAG and get its status
                            docs_by_status = await rag.get_docs_by_status("processed")
                            
                            # Look for our document by source identifier
                            for doc_id, doc_status in docs_by_status.items():
                                if doc_status.file_path == source_identifier:
                                    # Document is processed successfully
                                    db_inner = await get_namespace_data("document_requests")
                                    if request_id in db_inner:
                                        async with get_storage_lock():
                                            db_inner[request_id].update({
                                                "status": "Ready",
                                                "updatedAt": get_current_timestamp(),
                                                "downloadUrl": f"/webhook/api/document-requests/{request_id}/download"
                                            })
                                    return
                            
                            # Check for failed documents
                            failed_docs = await rag.get_docs_by_status("failed")
                            for doc_id, doc_status in failed_docs.items():
                                if doc_status.file_path == source_identifier:
                                    # Document processing failed
                                    db_inner = await get_namespace_data("document_requests")
                                    if request_id in db_inner:
                                        async with get_storage_lock():
                                            db_inner[request_id].update({
                                                "status": "Failed",
                                                "errorMessage": doc_status.error or "Processing failed",
                                                "updatedAt": get_current_timestamp()
                                            })
                                    return
                            
                            # Still processing, wait and try again
                            await asyncio.sleep(5)
                            attempts += 1
                            
                        except Exception as e:
                            logger.warning(f"Error monitoring LightRAG processing for {request_id}: {str(e)}")
                            await asyncio.sleep(5)
                            attempts += 1
                    
                    # Timeout - mark as failed
                    db_inner = await get_namespace_data("document_requests")
                    if request_id in db_inner:
                        async with get_storage_lock():
                            db_inner[request_id].update({
                                "status": "Failed",
                                "errorMessage": "Processing timeout after 5 minutes",
                                "updatedAt": get_current_timestamp()
                            })
                
                background_tasks.add_task(monitor_lightrag_processing)
                
            except Exception as e:
                logger.warning(f"Failed to create document request entry: {str(e)}")

            # Add to background tasks for processing
            background_tasks.add_task(
                pipeline_index_texts,
                rag,
                [processed_content],
                file_sources=[source_identifier],
            )

                return InsertResponse(
                    status="success",
                    message=f"Content from n8n webhook ({source_identifier}) successfully received. Processing will continue in background.",
                )
            except Exception as e:
                logger.error(f"Error /documents/n8n_webhook: {str(e)}")
                logger.error(traceback.format_exc())
                raise HTTPException(status_code=500, detail=str(e))

    @router.delete(
        "", response_model=ClearDocumentsResponse, dependencies=[Depends(combined_auth)]
    )
    async def clear_documents():
        """
        Clear all documents from the RAG system.

        This endpoint deletes all documents, entities, relationships, and files from the system.
        It uses the storage drop methods to properly clean up all data and removes all files
        from the input directory.

        Returns:
            ClearDocumentsResponse: A response object containing the status and message.
                - status="success":           All documents and files were successfully cleared.
                - status="partial_success":   Document clear job exit with some errors.
                - status="busy":              Operation could not be completed because the pipeline is busy.
                - status="fail":              All storage drop operations failed, with message
                - message: Detailed information about the operation results, including counts
                  of deleted files and any errors encountered.

        Raises:
            HTTPException: Raised when a serious error occurs during the clearing process,
                          with status code 500 and error details in the detail field.
        """
        from lightrag.kg.shared_storage import (
            get_namespace_data,
            get_pipeline_status_lock,
        )

        # Get pipeline status and lock
        pipeline_status = await get_namespace_data("pipeline_status")
        pipeline_status_lock = get_pipeline_status_lock()

        # Check and set status with lock
        async with pipeline_status_lock:
            if pipeline_status.get("busy", False):
                return ClearDocumentsResponse(
                    status="busy",
                    message="Cannot clear documents while pipeline is busy",
                )
            # Set busy to true
            pipeline_status.update(
                {
                    "busy": True,
                    "job_name": "Clearing Documents",
                    "job_start": datetime.now().isoformat(),
                    "docs": 0,
                    "batchs": 0,
                    "cur_batch": 0,
                    "request_pending": False,  # Clear any previous request
                    "latest_message": "Starting document clearing process",
                }
            )
            # Cleaning history_messages without breaking it as a shared list object
            del pipeline_status["history_messages"][:]
            pipeline_status["history_messages"].append(
                "Starting document clearing process"
            )

        try:
            # Use drop method to clear all data
            drop_tasks = []
            storages = [
                rag.text_chunks,
                rag.full_docs,
                rag.entities_vdb,
                rag.relationships_vdb,
                rag.chunks_vdb,
                rag.chunk_entity_relation_graph,
                rag.doc_status,
            ]

            # Log storage drop start
            if "history_messages" in pipeline_status:
                pipeline_status["history_messages"].append(
                    "Starting to drop storage components"
                )

            for storage in storages:
                if storage is not None:
                    drop_tasks.append(storage.drop())

            # Wait for all drop tasks to complete
            drop_results = await asyncio.gather(*drop_tasks, return_exceptions=True)

            # Check for errors and log results
            errors = []
            storage_success_count = 0
            storage_error_count = 0

            for i, result in enumerate(drop_results):
                storage_name = storages[i].__class__.__name__
                if isinstance(result, Exception):
                    error_msg = f"Error dropping {storage_name}: {str(result)}"
                    errors.append(error_msg)
                    logger.error(error_msg)
                    storage_error_count += 1
                else:
                    logger.info(f"Successfully dropped {storage_name}")
                    storage_success_count += 1

            # Log storage drop results
            if "history_messages" in pipeline_status:
                if storage_error_count > 0:
                    pipeline_status["history_messages"].append(
                        f"Dropped {storage_success_count} storage components with {storage_error_count} errors"
                    )
                else:
                    pipeline_status["history_messages"].append(
                        f"Successfully dropped all {storage_success_count} storage components"
                    )

            # If all storage operations failed, return error status and don't proceed with file deletion
            if storage_success_count == 0 and storage_error_count > 0:
                error_message = "All storage drop operations failed. Aborting document clearing process."
                logger.error(error_message)
                if "history_messages" in pipeline_status:
                    pipeline_status["history_messages"].append(error_message)
                return ClearDocumentsResponse(status="fail", message=error_message)

            # Log file deletion start
            if "history_messages" in pipeline_status:
                pipeline_status["history_messages"].append(
                    "Starting to delete files in input directory"
                )

            # Delete all files in input_dir
            deleted_files_count = 0
            file_errors_count = 0

            for file_path in doc_manager.input_dir.glob("**/*"):
                if file_path.is_file():
                    try:
                        file_path.unlink()
                        deleted_files_count += 1
                    except Exception as e:
                        logger.error(f"Error deleting file {file_path}: {str(e)}")
                        file_errors_count += 1

            # Log file deletion results
            if "history_messages" in pipeline_status:
                if file_errors_count > 0:
                    pipeline_status["history_messages"].append(
                        f"Deleted {deleted_files_count} files with {file_errors_count} errors"
                    )
                    errors.append(f"Failed to delete {file_errors_count} files")
                else:
                    pipeline_status["history_messages"].append(
                        f"Successfully deleted {deleted_files_count} files"
                    )

            # Prepare final result message
            final_message = ""
            if errors:
                final_message = f"Cleared documents with some errors. Deleted {deleted_files_count} files."
                status = "partial_success"
            else:
                final_message = f"All documents cleared successfully. Deleted {deleted_files_count} files."
                status = "success"

            # Log final result
            if "history_messages" in pipeline_status:
                pipeline_status["history_messages"].append(final_message)

            # Return response based on results
            return ClearDocumentsResponse(status=status, message=final_message)
        except Exception as e:
            error_msg = f"Error clearing documents: {str(e)}"
            logger.error(error_msg)
            logger.error(traceback.format_exc())
            if "history_messages" in pipeline_status:
                pipeline_status["history_messages"].append(error_msg)
            raise HTTPException(status_code=500, detail=str(e))
        finally:
            # Reset busy status after completion
            async with pipeline_status_lock:
                pipeline_status["busy"] = False
                completion_msg = "Document clearing process completed"
                pipeline_status["latest_message"] = completion_msg
                if "history_messages" in pipeline_status:
                    pipeline_status["history_messages"].append(completion_msg)

    @router.get(
        "/pipeline_status",
        dependencies=[Depends(combined_auth)],
        response_model=PipelineStatusResponse,
    )
    async def get_pipeline_status() -> PipelineStatusResponse:
        """
        Get the current status of the document indexing pipeline.

        This endpoint returns information about the current state of the document processing pipeline,
        including the processing status, progress information, and history messages.

        Returns:
            PipelineStatusResponse: A response object containing:
                - autoscanned (bool): Whether auto-scan has started
                - busy (bool): Whether the pipeline is currently busy
                - job_name (str): Current job name (e.g., indexing files/indexing texts)
                - job_start (str, optional): Job start time as ISO format string
                - docs (int): Total number of documents to be indexed
                - batchs (int): Number of batches for processing documents
                - cur_batch (int): Current processing batch
                - request_pending (bool): Flag for pending request for processing
                - latest_message (str): Latest message from pipeline processing
                - history_messages (List[str], optional): List of history messages

        Raises:
            HTTPException: If an error occurs while retrieving pipeline status (500)
        """
        try:
            from lightrag.kg.shared_storage import (
                get_namespace_data,
                get_all_update_flags_status,
            )
            # Caching layer to throttle heavy status calculations
            import os, time
            CACHE_TTL = int(os.getenv("FRONTEND_POLL_TTL_SEC", "20"))
            if not hasattr(get_pipeline_status, "_cache"):
                get_pipeline_status._cache = {"ts": 0.0, "data": None}  # type: ignore[attr-defined]
            now = time.time()
            cached = get_pipeline_status._cache  # type: ignore[attr-defined]
            if cached["data"] is not None and (now - cached["ts"]) < CACHE_TTL:
                return cached["data"]

            pipeline_status = await get_namespace_data("pipeline_status")

            # Get update flags status for all namespaces
            update_status = await get_all_update_flags_status()

            # Convert MutableBoolean objects to regular boolean values
            processed_update_status = {}
            for namespace, flags in update_status.items():
                processed_flags = []
                for flag in flags:
                    # Handle both multiprocess and single process cases
                    if hasattr(flag, "value"):
                        processed_flags.append(bool(flag.value))
                    else:
                        processed_flags.append(bool(flag))
                processed_update_status[namespace] = processed_flags

            # Convert to regular dict if it's a Manager.dict
            status_dict = dict(pipeline_status)

            # Add processed update_status to the status dictionary
            status_dict["update_status"] = processed_update_status

            # Convert history_messages to a regular list if it's a Manager.list
            if "history_messages" in status_dict:
                status_dict["history_messages"] = list(status_dict["history_messages"])

            # Ensure job_start is properly formatted as a string with timezone information
            if "job_start" in status_dict and status_dict["job_start"]:
                # Use format_datetime to ensure consistent formatting
                status_dict["job_start"] = format_datetime(status_dict["job_start"])

            resp = PipelineStatusResponse(**status_dict)
            cached["data"] = resp
            cached["ts"] = now
            return resp
        except Exception as e:
            logger.error(f"Error getting pipeline status: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.get(
        "", response_model=DocsStatusesResponse, dependencies=[Depends(combined_auth)]
    )
    async def documents() -> DocsStatusesResponse:
        """
        Get the status of all documents in the system.

        This endpoint retrieves the current status of all documents, grouped by their
        processing status (PENDING, PROCESSING, PROCESSED, FAILED).

        Returns:
            DocsStatusesResponse: A response object containing a dictionary where keys are
                                DocStatus values and values are lists of DocStatusResponse
                                objects representing documents in each status category.

        Raises:
            HTTPException: If an error occurs while retrieving document statuses (500).
        """
        try:
            # Caching layer for frequent polling
            import os, time
            CACHE_TTL = int(os.getenv("FRONTEND_POLL_TTL_SEC", "20"))
            if not hasattr(documents, "_cache"):
                documents._cache = {"ts": 0.0, "data": None}  # type: ignore[attr-defined]
            now = time.time()
            cached = documents._cache  # type: ignore[attr-defined]
            if cached["data"] is not None and (now - cached["ts"]) < CACHE_TTL:
                return cached["data"]

            statuses = (
                DocStatus.PENDING,
                DocStatus.PROCESSING,
                DocStatus.PROCESSED,
                DocStatus.FAILED,
            )

            tasks = [rag.get_docs_by_status(status) for status in statuses]
            results: List[Dict[str, DocProcessingStatus]] = await asyncio.gather(*tasks)

            response = DocsStatusesResponse()

            for idx, result in enumerate(results):
                status = statuses[idx]
                for doc_id, doc_status in result.items():
                    if status not in response.statuses:
                        response.statuses[status] = []
                    response.statuses[status].append(
                        DocStatusResponse(
                            id=doc_id,
                            content_summary=doc_status.content_summary,
                            content_length=doc_status.content_length,
                            status=doc_status.status,
                            created_at=format_datetime(doc_status.created_at),
                            updated_at=format_datetime(doc_status.updated_at),
                            chunks_count=doc_status.chunks_count,
                            error=doc_status.error,
                            metadata=doc_status.metadata,
                            file_path=doc_status.file_path,
                        )
                    )

            cached_resp = DocsStatusesResponse(statuses=response.statuses)
            cached["data"] = cached_resp
            cached["ts"] = now
            return cached_resp
        except Exception as e:
            logger.error(f"Error getting documents status: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/clear_cache",
        response_model=ClearCacheResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def clear_cache(request: ClearCacheRequest):
        """
        Clear cache data from the LLM response cache storage.

        This endpoint allows clearing specific modes of cache or all cache if no modes are specified.
        Valid modes include: "default", "naive", "local", "global", "hybrid", "mix".
        - "default" represents extraction cache.
        - Other modes correspond to different query modes.

        Args:
            request (ClearCacheRequest): The request body containing optional modes to clear.

        Returns:
            ClearCacheResponse: A response object containing the status and message.

        Raises:
            HTTPException: If an error occurs during cache clearing (400 for invalid modes, 500 for other errors).
        """
        try:
            # Validate modes if provided
            valid_modes = ["default", "naive", "local", "global", "hybrid", "mix"]
            if request.modes and not all(mode in valid_modes for mode in request.modes):
                invalid_modes = [
                    mode for mode in request.modes if mode not in valid_modes
                ]
                raise HTTPException(
                    status_code=400,
                    detail=f"Invalid mode(s): {invalid_modes}. Valid modes are: {valid_modes}",
                )

            # Call the aclear_cache method
            await rag.aclear_cache(request.modes)

            # Prepare success message
            if request.modes:
                message = f"Successfully cleared cache for modes: {request.modes}"
            else:
                message = "Successfully cleared all cache"

            return ClearCacheResponse(status="success", message=message)
        except HTTPException:
            # Re-raise HTTP exceptions
            raise
        except Exception as e:
            logger.error(f"Error clearing cache: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    return router
